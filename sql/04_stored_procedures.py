"""deploy_pptx_docx_sp.py - XLSX実績方式と完全同一パターンの PPTX/Word SP"""
import os
import snowflake.connector

conn = snowflake.connector.connect(connection_name=os.getenv("SNOWFLAKE_CONNECTION_NAME", "KMOT_DEMO1"))
cur = conn.cursor()
cur.execute("USE ROLE ACCOUNTADMIN")
cur.execute("USE DATABASE NIPPONLIFE_DEMO_DB")
cur.execute("USE SCHEMA RAW")
cur.execute("USE WAREHOUSE NIPPONLIFE_DEMO_WH")

cur.execute("CREATE STAGE IF NOT EXISTS NIPPONLIFE_DEMO_DB.RAW.PROPOSAL_EXPORT_STAGE")
print("Stage created")

# ============================================================
# PPTX SP - XLSX方式を完全コピー
# ============================================================
pptx_sql = """
CREATE OR REPLACE PROCEDURE NIPPONLIFE_DEMO_DB.RAW.GENERATE_PROPOSAL_PPTX(
    COMPANY_NAME VARCHAR,
    PRODUCT_NAMES VARCHAR,
    PROPOSAL_CONTENT VARCHAR
)
RETURNS VARCHAR
LANGUAGE PYTHON
RUNTIME_VERSION = '3.11'
PACKAGES = ('snowflake-snowpark-python', 'python-pptx')
HANDLER = 'main'
EXECUTE AS CALLER
AS
$$
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from datetime import datetime
import os
import re

def main(session, company_name: str, product_names: str, proposal_content: str) -> str:
    RED = RGBColor(0xE6, 0x00, 0x12)
    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)

    def add_tb(slide, l, t, w, h, text, sz=12, bold=False, color=None):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        r = p.runs[0]
        r.font.size = Pt(sz)
        r.font.bold = bold
        if color:
            r.font.color.rgb = color
        return tf

    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = RED
    hdr.line.fill.background()
    add_tb(s1, 0.5, 1.2, 12, 1.0, company_name + ' 御中', 26, True)
    add_tb(s1, 0.5, 2.5, 12, 0.7, '退職給付・福利厚生制度の最適化に向けたご提案', 16)
    today_str = datetime.now().strftime('%Y年%m月%d日')
    add_tb(s1, 0.5, 3.3, 12, 0.5, '日本生命保険相互会社 法人部 | ' + today_str, 11)
    add_tb(s1, 0.5, 3.9, 12, 0.4, '提案商品: ' + product_names, 10)

    NL2 = chr(10) + chr(10)
    NL  = chr(10)
    sections = [s.strip() for s in proposal_content.split(NL2) if s.strip()]
    for sec in sections[:6]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        hdr2 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.7))
        hdr2.fill.solid()
        hdr2.fill.fore_color.rgb = RED
        hdr2.line.fill.background()
        lines = sec.strip().split(NL)
        title = lines[0] if lines else ''
        body_text = NL.join(lines[1:]) if len(lines) > 1 else ''
        add_tb(slide, 0.2, 0.1, 12.5, 0.5, title, 14, True, RGBColor(0xFF, 0xFF, 0xFF))
        if body_text:
            add_tb(slide, 0.5, 0.9, 12, 4, body_text, 11)

    safe_name = re.sub(r'[^\\w\\-.]', '_', company_name + '_' + datetime.now().strftime('%Y%m%d'))
    if not safe_name.endswith('.pptx'):
        safe_name = safe_name + '.pptx'

    local_path = '/tmp/' + safe_name
    prs.save(local_path)

    stage_path = '@NIPPONLIFE_DEMO_DB.RAW.PROPOSAL_EXPORT_STAGE'
    session.file.put(local_path, stage_path, auto_compress=False, overwrite=True)

    os.remove(local_path)

    stage_files = session.sql(f"LIST {stage_path} PATTERN='.*{safe_name}.*'").collect()
    actual_name = safe_name
    for f in stage_files:
        fname = f['name'].split('/')[-1]
        if safe_name in fname:
            actual_name = fname
            break

    presigned_url_result = session.sql(
        f"SELECT GET_PRESIGNED_URL({stage_path}, '{actual_name}', 3600) AS URL"
    ).collect()
    download_url = presigned_url_result[0]['URL']

    return f'[{company_name}_提案書.pptxをダウンロード]({download_url})'
$$;
"""

# ============================================================
# Word SP - XLSX方式を完全コピー
# ============================================================
docx_sql = """
CREATE OR REPLACE PROCEDURE NIPPONLIFE_DEMO_DB.RAW.GENERATE_PROPOSAL_DOCX(
    COMPANY_NAME VARCHAR,
    PRODUCT_NAMES VARCHAR,
    PROPOSAL_CONTENT VARCHAR
)
RETURNS VARCHAR
LANGUAGE PYTHON
RUNTIME_VERSION = '3.11'
PACKAGES = ('snowflake-snowpark-python', 'python-docx')
HANDLER = 'main'
EXECUTE AS CALLER
AS
$$
from docx import Document
from docx.shared import Pt, RGBColor
from datetime import datetime
import os
import re

def main(session, company_name: str, product_names: str, proposal_content: str) -> str:
    doc = Document()
    NL2 = chr(10) + chr(10)
    NL  = chr(10)

    h0 = doc.add_paragraph()
    run0 = h0.add_run(company_name + '  御中')
    run0.font.size = Pt(20)
    run0.font.bold = True
    run0.font.color.rgb = RGBColor(0xE6, 0x00, 0x12)

    h1 = doc.add_paragraph()
    run1 = h1.add_run('退職給付・福利厚生制度の最適化に向けたご提案')
    run1.font.size = Pt(14)
    run1.font.bold = True

    today_str = datetime.now().strftime('%Y年%m月%d日')
    doc.add_paragraph('提案日: ' + today_str)
    doc.add_paragraph('提案商品: ' + product_names)
    doc.add_paragraph('日本生命保険相互会社 法人部')
    doc.add_paragraph('')

    sections = [s.strip() for s in proposal_content.split(NL2) if s.strip()]
    for sec in sections:
        lines = sec.strip().split(NL)
        if not lines:
            continue
        title_p = doc.add_paragraph()
        title_run = title_p.add_run(lines[0])
        title_run.font.size = Pt(13)
        title_run.font.bold = True
        for line in lines[1:]:
            stripped = line.strip()
            if stripped:
                body_p = doc.add_paragraph()
                body_run = body_p.add_run(stripped)
                body_run.font.size = Pt(11)
        doc.add_paragraph('')

    safe_name = re.sub(r'[^\\w\\-.]', '_', company_name + '_' + datetime.now().strftime('%Y%m%d'))
    if not safe_name.endswith('.docx'):
        safe_name = safe_name + '.docx'

    local_path = '/tmp/' + safe_name
    doc.save(local_path)

    stage_path = '@NIPPONLIFE_DEMO_DB.RAW.PROPOSAL_EXPORT_STAGE'
    session.file.put(local_path, stage_path, auto_compress=False, overwrite=True)

    os.remove(local_path)

    stage_files = session.sql(f"LIST {stage_path} PATTERN='.*{safe_name}.*'").collect()
    actual_name = safe_name
    for f in stage_files:
        fname = f['name'].split('/')[-1]
        if safe_name in fname:
            actual_name = fname
            break

    presigned_url_result = session.sql(
        f"SELECT GET_PRESIGNED_URL({stage_path}, '{actual_name}', 3600) AS URL"
    ).collect()
    download_url = presigned_url_result[0]['URL']

    return f'[{company_name}_提案書.docxをダウンロード]({download_url})'
$$;
"""

for name, ddl in [("PPTX", pptx_sql), ("DOCX", docx_sql)]:
    try:
        cur.execute(ddl)
        print(f"Created: GENERATE_PROPOSAL_{name}")
    except Exception as e:
        print(f"Error {name}: {e}")

# テスト実行
print("\n--- Testing PPTX ---")
try:
    cur.execute("""CALL NIPPONLIFE_DEMO_DB.RAW.GENERATE_PROPOSAL_PPTX(
        'KDDI', 'DC, GLTD',
        '1. ご提案の背景

KDDIは大規模採用を進めています。

2. 課題

・福利厚生の充実が必要'
    )""")
    result = cur.fetchone()[0]
    print("PPTX result:", result[:200])
except Exception as e:
    print(f"PPTX test error: {e}")

print("\n--- Testing DOCX ---")
try:
    cur.execute("""CALL NIPPONLIFE_DEMO_DB.RAW.GENERATE_PROPOSAL_DOCX(
        'KDDI', 'DC, GLTD',
        '1. ご提案の背景

KDDIは大規模採用を進めています。

2. 課題

・福利厚生の充実が必要'
    )""")
    result = cur.fetchone()[0]
    print("DOCX result:", result[:200])
except Exception as e:
    print(f"DOCX test error: {e}")

cur.close()
conn.close()
print("\nDone!")
