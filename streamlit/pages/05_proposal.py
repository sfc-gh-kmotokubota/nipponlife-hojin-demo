"""
05_📄_提案書.py - 提案書自動生成（PPTX + Word） (F-05)
"""
import streamlit as st
import io
from datetime import date
from snowflake.snowpark.context import get_active_session

st.set_page_config(page_title="提案書自動生成", layout="wide")
st.title("📄 提案書自動生成（PowerPoint / Word）")
st.caption("企業名と商品を選んでボタン1つ。最新情報が差し込まれた提案書が30秒で完成。既存テンプレートにも対応。")

session = get_active_session()

# ────────────────────────────────────────
# 入力フォーム
# ────────────────────────────────────────
col1, col2 = st.columns(2)

with col1:
    companies = session.sql("""
        SELECT COMPANY_ID, COMPANY_NAME, EMPLOYEE_COUNT, INDUSTRY_LARGE, PENSION_TYPE
        FROM NIPPONLIFE_DEMO_DB.RAW.T_CUSTOMER_COMPANIES ORDER BY COMPANY_NAME
    """).to_pandas()
    company_options = {r["COMPANY_NAME"]: r["COMPANY_ID"] for _, r in companies.iterrows()}
    company_names = list(company_options.keys())

    # alertページからの遷移時に会社名を自動選択
    preselect_name = None
    if "proposal_company_name" in st.session_state:
        preselect_name = st.session_state.pop("proposal_company_name")
    elif "proposal_company" in st.session_state:
        cid = st.session_state.pop("proposal_company")
        for n, c in company_options.items():
            if c == cid:
                preselect_name = n
                break

    default_idx = company_names.index(preselect_name) if preselect_name in company_names else 0
    selected_name = st.selectbox("対象企業", company_names, index=default_idx)
    selected_cid = company_options[selected_name]

    proposal_date = st.date_input("提案日", value=date.today())

    contacts = session.sql(f"""
        SELECT FULL_NAME, TITLE FROM NIPPONLIFE_DEMO_DB.RAW.T_CONTACTS
        WHERE COMPANY_ID = '{selected_cid}'
        ORDER BY IS_PRIMARY DESC
    """).to_pandas()
    if not contacts.empty:
        contact_options = [f"{r['FULL_NAME']}（{r['TITLE']}）" for _, r in contacts.iterrows()]
        recipient = st.selectbox("提出先担当者", contact_options)
    else:
        recipient = st.text_input("提出先担当者", "人事部長 様")

with col2:
    products = session.sql("""
        SELECT PRODUCT_ID, PRODUCT_NAME, PRODUCT_CATEGORY
        FROM NIPPONLIFE_DEMO_DB.RAW.T_INSURANCE_PRODUCTS
        ORDER BY PRODUCT_CATEGORY, PRODUCT_ID
    """).to_pandas()

    selected_products = st.multiselect(
        "提案する商品（複数選択可）",
        products["PRODUCT_NAME"].tolist(),
        default=products["PRODUCT_NAME"].head(2).tolist()
    )

    proposal_style = st.radio("提案書スタイル", ["標準版（12ページ）", "簡易版（5ページ）"])

    custom_template = st.file_uploader(
        "📁 既存テンプレートをアップロード（オプション）",
        type=["pptx","docx"],
        help="既存の提案書テンプレートがある場合はアップロードしてください"
    )

# 自動差し込み情報表示
company_data = companies[companies["COMPANY_ID"] == selected_cid].iloc[0]
alerts = session.sql(f"""
    SELECT EVENT_TYPE, EVENT_SUMMARY FROM NIPPONLIFE_DEMO_DB.RAW.T_EVENT_ALERTS
    WHERE COMPANY_ID = '{selected_cid}' AND STATUS = 'UNREAD'
    ORDER BY CASE INSURANCE_RELEVANCE WHEN '最高' THEN 1 WHEN '高' THEN 2 ELSE 3 END LIMIT 2
""").to_pandas()

with st.expander("📊 自動差し込み情報（AIが企業データから生成）", expanded=True):
    ai_cols = st.columns(3)
    with ai_cols[0]:
        st.markdown("**企業現状**")
        st.markdown(f"✅ 従業員数: {company_data['EMPLOYEE_COUNT']:,}名")
        st.markdown(f"✅ 業種: {company_data['INDUSTRY_LARGE']}")
        st.markdown(f"✅ 年金制度: {company_data['PENSION_TYPE'] or '未確認'}")
    with ai_cols[1]:
        st.markdown("**最新事業イベント**")
        if not alerts.empty:
            for _, a in alerts.iterrows():
                st.markdown(f"✅ {a['EVENT_TYPE']}: {str(a['EVENT_SUMMARY'])[:50]}...")
        else:
            st.markdown("✅ 直近ニュースから要約済み")
    with ai_cols[2]:
        st.markdown("**市場環境**")
        st.markdown("✅ 10年金利: 1.45%（上昇局面）")
        st.markdown("✅ DC移行提案の最適タイミング")
        st.markdown("✅ 金利シグナル: 積立改善傾向")

# ────────────────────────────────────────
# 生成ボタン
# ────────────────────────────────────────
if st.button("▶ AI 提案書ドラフト生成（約30秒）", type="primary", disabled=not selected_products):
    with st.spinner("AI が提案書を生成中..."):
        products_text = "、".join(selected_products[:3])

        # AI でコンテンツ生成
        try:
            content_raw = session.sql(f"""
                SELECT AI_COMPLETE('mistral-large2', CONCAT(
                    '以下の情報を元に、日本生命保険の法人向け提案書の主要セクションを生成してください。
                     条件: マークダウン記法（#、##、**、*）は一切使わないこと。セクション見出しは「1. XXX」形式のみ。箇条書きは「・」のみ使用。
                     企業名: {selected_name}
                     業種: {company_data["INDUSTRY_LARGE"]}
                     従業員数: {company_data["EMPLOYEE_COUNT"]:,}名
                     提案商品: {products_text}
                     最新事業イベント: {alerts.iloc[0]["EVENT_TYPE"] if not alerts.empty else "なし"}
                     市場環境: 10年金利1.45%（上昇局面）
                     以下のセクションを日本語で生成:
                     1. ご提案の背景（2-3文）
                     2. 御社の現状課題（3点箇条書き）
                     3. 日本生命のご提案内容（各商品2-3文）
                     4. 期待される効果（3点）
                     5. 今後のスケジュール案（3ステップ）')) AS CONTENT
            """).collect()[0]["CONTENT"]
            import re
            content = re.sub(r'^#{1,6}\s*', '', content_raw, flags=re.MULTILINE)
            content = re.sub(r'\*\*(.+?)\*\*', r'\1', content)
            content = re.sub(r'\*(.+?)\*', r'\1', content)
            content = content.replace('\\n', '\n').replace('\\t', ' ')
            content = content.strip('"').strip("'").strip()
        except:
            content = f"""1. ご提案の背景
{selected_name}様の従業員{company_data['EMPLOYEE_COUNT']:,}名を守る保障制度の充実を目的に、この度ご提案させていただきます。足元の金利環境（10年金利1.45%）は制度見直しの最適タイミングを迎えております。

2. 御社の現状課題
・現行の退職給付制度における将来的な積立リスクの管理
・従業員の長期就業不能リスクへの対応（GLTD未整備）
・健康経営推進に伴う福利厚生制度の充実

3. 日本生命のご提案内容
{products_text}を組み合わせた包括的な保障スキームをご提案いたします。

4. 期待される効果
・従業員エンゲージメントの向上と採用競争力の強化
・企業の退職給付費用の安定化とリスク低減
・健康経営優良法人認定取得への貢献

5. 今後のスケジュール案
STEP1: 詳細ヒアリング（来月第1週） → STEP2: 制度設計案の提示（来月末） → STEP3: 最終合意（再来月）"""

        st.session_state["proposal_content"] = content
        st.session_state["proposal_company"] = selected_name

if "proposal_content" in st.session_state:
    st.markdown("---")
    st.subheader("📋 生成プレビュー")
    st.text_area("提案書本文（編集可）", value=st.session_state["proposal_content"], height=300, key="proposal_edit")

    dl_col1, dl_col2 = st.columns(2)

    with dl_col1:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor

            prs = Presentation()
            prs.slide_width = Emu(9144000)
            prs.slide_height = Emu(5143500)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(2))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{selected_name} 御中"
            run = p.runs[0]
            run.font.size = Pt(20)

            p2 = tf.add_paragraph()
            p2.text = f"退職給付制度最適化のご提案 - {proposal_date}"
            run2 = p2.runs[0]
            run2.font.size = Pt(16)

            p3 = tf.add_paragraph()
            p3.text = "日本生命保険相互会社 法人部"
            run3 = p3.runs[0]
            run3.font.size = Pt(14)

            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(4.5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p_head = tf2.paragraphs[0]
            p_head.text = "ご提案内容"
            run_h = p_head.runs[0]
            run_h.font.size = Pt(18)
            run_h.font.bold = True
            run_h.font.color.rgb = RGBColor(0xE6, 0x00, 0x12)

            for line in st.session_state["proposal_content"].split("\n")[:15]:
                if line.strip():
                    p_new = tf2.add_paragraph()
                    p_new.text = line
                    r = p_new.runs[0]
                    r.font.size = Pt(11)

            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_data = pptx_buffer.getvalue()

            st.download_button(
                "📊 PowerPoint ⬇ ダウンロード",
                data=pptx_data,
                file_name=f"nissay_proposal_{selected_name}_{proposal_date}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except ImportError:
            st.warning("python-pptx が必要です")
        except Exception as e:
            st.error(f"PPTX生成エラー: {e}")

    with dl_col2:
        # Word生成
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor as DocxRGB
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()
            doc.add_heading(f"{selected_name} 御中", 0)
            doc.add_heading("退職給付制度最適化のご提案", 1)
            doc.add_paragraph(f"提案日: {proposal_date} | 提出: 日本生命保険相互会社 法人部")
            doc.add_paragraph()

            for line in st.session_state["proposal_content"].split("\n"):
                if line.strip():
                    if line.strip()[0].isdigit() and "." in line:
                        doc.add_heading(line.strip(), 2)
                    else:
                        para = doc.add_paragraph(line.strip())
                        para.style.font.size = Pt(11)

            docx_buffer = io.BytesIO()
            doc.save(docx_buffer)
            docx_buffer.seek(0)

            st.download_button(
                "📄 Word ⬇ ダウンロード",
                data=docx_buffer.getvalue(),
                file_name=f"nissay_proposal_{selected_name}_{proposal_date}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
        except ImportError:
            st.warning("python-docx が必要です")
        except Exception as e:
            st.error(f"Word生成エラー: {e}")
